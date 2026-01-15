#!/usr/bin/env python3
"""
Generate PowerPoint presentation matching the dark HTML design.
Dark theme with neon yellow accents, bold typography.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

def set_slide_background(slide, rgb_color):
    """Set solid color background for a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_color)

def add_text_box(slide, left, top, width, height, text, font_size, bold=False, 
                 color=(255, 255, 255), alignment=PP_ALIGN.LEFT):
    """Helper to add text box with styling."""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.alignment = alignment
    
    return textbox

def add_rounded_rectangle(slide, left, top, width, height, fill_color, text_content=None, 
                         border_color=None, border_width=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    
    # Set fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_color)
    
    # Set border
    if border_color:
        shape.line.color.rgb = RGBColor(*border_color)
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    
    # Add text if provided
    if text_content:
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = text_content
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    return shape

def create_slide_1_problem(prs):
    """Slide 1: The Problem"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide, (0, 0, 0))
    
    # Header: "01. THE PROBLEM"
    add_text_box(slide, 0.5, 0.5, 6, 0.5, "01. THE PROBLEM", 36, bold=True, 
                 color=(204, 255, 0))
    
    # Top right: LOSS info
    add_text_box(slide, 7, 0.5, 2.5, 0.5, "LOSS: $200,000 / PROJECT", 10, 
                 color=(120, 120, 120), alignment=PP_ALIGN.RIGHT)
    
    # Main headline
    add_text_box(slide, 0.8, 2, 5, 1.5, "AI Storyboards are", 60, bold=True)
    add_text_box(slide, 0.8, 3.2, 5, 1, "Broken.", 60, bold=True, color=(204, 255, 0))
    
    # Subtext
    add_text_box(slide, 0.8, 4.5, 4.5, 0.8, 
                 'Characters change in every frame. This "Character Drift" makes AI unusable for professional film sets.', 
                 18, color=(180, 180, 180))
    
    # Stats: 60%
    add_text_box(slide, 0.8, 5.5, 1.5, 0.8, "60%", 48, bold=True)
    add_text_box(slide, 0.8, 6.2, 1.5, 0.3, "Time Wasted Fixing", 10, color=(120, 120, 120))
    
    # Stats: 30%
    add_text_box(slide, 2.5, 5.5, 1.5, 0.8, "30%", 48, bold=True)
    add_text_box(slide, 2.5, 6.2, 1.5, 0.3, "Error Rate", 10, color=(120, 120, 120))
    
    # Visual on right (Scene circles)
    add_rounded_rectangle(slide, 6, 2.5, 3.5, 3.5, (30, 30, 30))
    
    # Scene 1 circle
    scene1 = slide.shapes.add_shape(
        9,  # Oval
        Inches(6.8), Inches(3), Inches(1.2), Inches(1.2)
    )
    scene1.fill.solid()
    scene1.fill.fore_color.rgb = RGBColor(40, 40, 40)
    scene1.line.color.rgb = RGBColor(239, 68, 68)  # Red
    scene1.line.width = Pt(4)
    text_frame = scene1.text_frame
    p = text_frame.paragraphs[0]
    p.text = "SCENE 1"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = RGBColor(239, 68, 68)
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Scene 2 circle
    scene2 = slide.shapes.add_shape(
        9,  # Oval
        Inches(6.8), Inches(4.5), Inches(1.2), Inches(1.2)
    )
    scene2.fill.solid()
    scene2.fill.fore_color.rgb = RGBColor(40, 40, 40)
    scene2.line.color.rgb = RGBColor(251, 146, 60)  # Orange
    scene2.line.width = Pt(4)
    text_frame = scene2.text_frame
    p = text_frame.paragraphs[0]
    p.text = "? SCENE 2"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.italic = True
    p.font.color.rgb = RGBColor(251, 146, 60)
    p.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def create_slide_2_solution(prs):
    """Slide 2: The Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, (0, 0, 0))
    
    # Header
    add_text_box(slide, 0.5, 0.5, 6, 0.5, "02. THE SOLUTION", 36, bold=True, 
                 color=(204, 255, 0))
    
    # Main headline
    add_text_box(slide, 0.8, 1.8, 5, 0.8, "Meet", 60, bold=True)
    add_text_box(slide, 2.5, 1.8, 5, 0.8, "CharacterLock.", 60, bold=True, 
                 color=(204, 255, 0))
    
    # Feature cards
    card_y = 3.2
    card_spacing = 3.2
    
    # Card 1: Digital Identity
    card1 = add_rounded_rectangle(slide, 0.8, card_y, 2.8, 2.5, (30, 30, 30), 
                                  border_color=(60, 60, 60), border_width=1)
    # Number badge
    badge1 = add_rounded_rectangle(slide, 1, card_y + 0.3, 0.4, 0.4, (204, 255, 0))
    badge1.text_frame.paragraphs[0].text = "01"
    badge1.text_frame.paragraphs[0].font.size = Pt(14)
    badge1.text_frame.paragraphs[0].font.bold = True
    badge1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    # Text
    add_text_box(slide, 1, card_y + 0.9, 2.5, 0.4, "Digital Identity", 20, bold=True)
    add_text_box(slide, 1, card_y + 1.4, 2.5, 0.8, 
                 "Lock the character's facial DNA so they never change.", 12, 
                 color=(180, 180, 180))
    
    # Card 2: Auto-Checker
    card2 = add_rounded_rectangle(slide, 3.8, card_y, 2.8, 2.5, (30, 30, 30), 
                                  border_color=(60, 60, 60), border_width=1)
    badge2 = add_rounded_rectangle(slide, 4.0, card_y + 0.3, 0.4, 0.4, (204, 255, 0))
    badge2.text_frame.paragraphs[0].text = "02"
    badge2.text_frame.paragraphs[0].font.size = Pt(14)
    badge2.text_frame.paragraphs[0].font.bold = True
    badge2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    add_text_box(slide, 4.0, card_y + 0.9, 2.5, 0.4, "Auto-Checker", 20, bold=True)
    add_text_box(slide, 4.0, card_y + 1.4, 2.5, 0.8, 
                 "Our AI scores every frame. If it's not perfect, we flag it.", 12, 
                 color=(180, 180, 180))
    
    # Card 3: Smart Repair
    card3 = add_rounded_rectangle(slide, 6.8, card_y, 2.8, 2.5, (30, 30, 30), 
                                  border_color=(60, 60, 60), border_width=1)
    badge3 = add_rounded_rectangle(slide, 7.0, card_y + 0.3, 0.4, 0.4, (204, 255, 0))
    badge3.text_frame.paragraphs[0].text = "03"
    badge3.text_frame.paragraphs[0].font.size = Pt(14)
    badge3.text_frame.paragraphs[0].font.bold = True
    badge3.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
    add_text_box(slide, 7.0, card_y + 0.9, 2.5, 0.4, "Smart Repair", 20, bold=True)
    add_text_box(slide, 7.0, card_y + 1.4, 2.5, 0.8, 
                 "One-click regeneration to fix inconsistencies instantly.", 12, 
                 color=(180, 180, 180))

def create_slide_3_impact(prs):
    """Slide 3: The Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, (0, 0, 0))
    
    # Header
    add_text_box(slide, 0.5, 0.5, 6, 0.5, "03. THE IMPACT", 36, bold=True, 
                 color=(204, 255, 0))
    
    # Main headline
    add_text_box(slide, 0.8, 1.8, 5, 0.8, "From Weeks to", 60, bold=True)
    add_text_box(slide, 0.8, 2.5, 5, 0.8, "Hours.", 60, bold=True, 
                 color=(204, 255, 0))
    
    # Impact points
    impact_y = 3.8
    
    # 01.
    add_text_box(slide, 0.8, impact_y, 0.5, 0.5, "01.", 48, bold=True, 
                 color=(60, 60, 60))
    add_text_box(slide, 1.5, impact_y + 0.1, 4, 0.4, 
                 "90% Reduction in storyboard costs.", 18, color=(200, 200, 200))
    
    # 02.
    add_text_box(slide, 0.8, impact_y + 0.7, 0.5, 0.5, "02.", 48, bold=True, 
                 color=(60, 60, 60))
    add_text_box(slide, 1.5, impact_y + 0.8, 4, 0.4, 
                 "Studio-quality continuity for indie budgets.", 18, color=(200, 200, 200))
    
    # 03.
    add_text_box(slide, 0.8, impact_y + 1.4, 0.5, 0.5, "03.", 48, bold=True, 
                 color=(60, 60, 60))
    add_text_box(slide, 1.5, impact_y + 1.5, 4, 0.4, 
                 "Tapping into a $2.8B global market.", 18, color=(200, 200, 200))
    
    # Chart container
    chart_container = add_rounded_rectangle(slide, 6, 3.5, 3.5, 3, (30, 30, 30), 
                                           border_color=(60, 60, 60), border_width=1)
    
    # Chart label
    add_text_box(slide, 6.3, 3.8, 3, 0.3, "EFFICIENCY GAIN", 9, bold=True, 
                 color=(120, 120, 120))
    
    # Bars
    # Traditional AI bar (tall, gray)
    add_rounded_rectangle(slide, 6.5, 4.5, 1, 1.5, (60, 60, 60))
    # CharacterLock bar (short, yellow)
    add_rounded_rectangle(slide, 7.7, 5.8, 1, 0.2, (204, 255, 0))
    
    # Labels
    add_text_box(slide, 6.3, 6.2, 1.3, 0.2, "TRADITIONAL AI", 8, bold=True, 
                 color=(120, 120, 120))
    add_text_box(slide, 7.7, 6.2, 1.5, 0.2, "CHARACTERLOCK", 8, bold=True, 
                 color=(120, 120, 120))

def create_slide_4_engine(prs):
    """Slide 4: The Engine"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, (0, 0, 0))
    
    # Header
    add_text_box(slide, 0.5, 0.5, 6, 0.5, "04. THE ENGINE", 36, bold=True, 
                 color=(204, 255, 0))
    
    # Main headline
    add_text_box(slide, 0.5, 1.8, 9, 0.8, "Simple. Scalable. Powerful.", 44, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    
    # Flow diagram
    flow_y = 3.5
    
    # Input box
    input_box = add_rounded_rectangle(slide, 1.5, flow_y, 2, 1.2, (0, 0, 0), 
                                     border_color=(60, 60, 60), border_width=2)
    input_box.line.dash_style = 2  # Dashed
    add_text_box(slide, 1.6, flow_y + 0.1, 1.8, 0.3, "INPUT", 10, 
                 color=(120, 120, 120), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1.6, flow_y + 0.5, 1.8, 0.5, "Script + Character", 16, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    
    # Arrow 1
    add_text_box(slide, 3.6, flow_y + 0.4, 0.5, 0.5, "→", 32, color=(204, 255, 0))
    
    # Engine box (yellow)
    engine_box = add_rounded_rectangle(slide, 4.2, flow_y, 2, 1.2, (204, 255, 0))
    add_text_box(slide, 4.3, flow_y + 0.1, 1.8, 0.3, "ENGINE", 10, 
                 color=(100, 100, 100), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 4.3, flow_y + 0.5, 1.8, 0.5, "CHARACTERLOCK", 16, bold=True, 
                 color=(0, 0, 0), alignment=PP_ALIGN.CENTER)
    
    # Arrow 2
    add_text_box(slide, 6.3, flow_y + 0.4, 0.5, 0.5, "→", 32, color=(204, 255, 0))
    
    # Output box
    output_box = add_rounded_rectangle(slide, 6.9, flow_y, 2, 1.2, (0, 0, 0), 
                                      border_color=(60, 60, 60), border_width=2)
    output_box.line.dash_style = 2  # Dashed
    add_text_box(slide, 7.0, flow_y + 0.1, 1.8, 0.3, "OUTPUT", 10, 
                 color=(120, 120, 120), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 7.0, flow_y + 0.5, 1.8, 0.5, "Consistent Storyboard", 16, bold=True, 
                 alignment=PP_ALIGN.CENTER)
    
    # Bottom text
    add_text_box(slide, 2, 5.5, 6.5, 0.8, 
                 "Built on top of OpenAI's Vision and Generation infrastructure, optimized with our proprietary \"Identity Scoring\" logic.", 
                 14, color=(120, 120, 120), alignment=PP_ALIGN.CENTER)

def main():
    """Generate the dark-themed presentation."""
    prs = Presentation()
    
    # Set slide size to match HTML (1280x720)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    print("Creating slide 1: The Problem...")
    create_slide_1_problem(prs)
    
    print("Creating slide 2: The Solution...")
    create_slide_2_solution(prs)
    
    print("Creating slide 3: The Impact...")
    create_slide_3_impact(prs)
    
    print("Creating slide 4: The Engine...")
    create_slide_4_engine(prs)
    
    # Save presentation
    output_file = "CharacterLock_Dark_Theme.pptx"
    prs.save(output_file)
    
    print(f"\n✓ Dark-themed presentation created: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  Theme: Black background with neon yellow accents")
    print(f"  Style: Matches HTML design")
    print(f"  Location: {output_file}")

if __name__ == "__main__":
    main()
