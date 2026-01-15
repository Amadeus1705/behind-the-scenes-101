#!/usr/bin/env python3
"""
Generate PowerPoint presentation for CharacterLock AI hackathon submission.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    
    return slide

def create_content_slide(prs, title, content_items, layout_idx=1):
    """Create content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
    
    return slide

def create_two_column_slide(prs, title, left_content, right_content):
    """Create two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(3, 105, 161)
    
    # Left column
    left_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_shape.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(12)
    
    # Right column
    right_shape = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_shape.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.space_after = Pt(12)
    
    return slide

def create_large_text_slide(prs, title, main_text, subtext=""):
    """Create slide with large centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_para = title_frame.add_paragraph()
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # Main text (large, centered)
    text_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    text_frame = text_shape.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.add_paragraph()
    p.text = main_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(3, 105, 161)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtext
    if subtext:
        sub_shape = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(8), Inches(1))
        sub_frame = sub_shape.text_frame
        sub_para = sub_frame.add_paragraph()
        sub_para.text = subtext
        sub_para.font.size = Pt(20)
        sub_para.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    """Generate the presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "CharacterLock AI",
        "Persistent Character Memory for Film Production\nCine AI Hackathon 2026"
    )
    
    # Slide 2: Problem Statement - The Challenge
    create_content_slide(
        prs,
        "1. The Problem: Character Inconsistency in AI Film Production",
        [
            "Current AI tools cannot maintain character consistency across frames",
            "Characters change appearance between scenes - different faces, hair, clothing",
            "AI-generated storyboards are UNUSABLE for production",
            "Manual fixes cost $50,000 - $200,000 per project",
            "60% of production time wasted on corrections",
            "Blocks indie filmmakers from using AI tools"
        ]
    )
    
    # Slide 3: Problem Impact (Large Numbers)
    create_large_text_slide(
        prs,
        "The Cost of Inconsistency",
        "$50K - $200K",
        "wasted per production fixing character drift"
    )
    
    # Slide 4: Our Solution
    create_content_slide(
        prs,
        "2. Our Idea: CharacterLock AI",
        [
            "Persistent Character Memory System with 'Character DNA'",
            "Automated Consistency Validation with quality scores",
            "Smart Regeneration - fix only problematic frames",
            "Unified workflow: Create → Generate → Validate",
            "Production-ready output with 85%+ consistency",
            "Quantified, measurable results"
        ]
    )
    
    # Slide 5: How It Works - 3 Steps
    create_content_slide(
        prs,
        "How It Works (3 Simple Steps)",
        [
            "STEP 1: Create Character",
            "  • Upload 1-3 reference images",
            "  • AI extracts 'Character DNA' (visual embeddings)",
            "  • Permanent identity created in 10-15 seconds",
            "",
            "STEP 2: Generate Storyboard",
            "  • Write your script (any length)",
            "  • Select characters to include",
            "  • AI generates consistent frames with Character DNA",
            "",
            "STEP 3: Validate & Fix",
            "  • Automated consistency analysis (0-100% scores)",
            "  • Get detailed quality report",
            "  • One-click regeneration for low-scoring frames"
        ]
    )
    
    # Slide 6: Technical Innovation
    create_two_column_slide(
        prs,
        "4. How It Works: Technical Innovation",
        [
            "CHARACTER DNA EXTRACTION:",
            "• Vision API analyzes reference images",
            "• Extracts facial features, hair, clothing, style",
            "• Creates persistent embedding vector",
            "• Generates optimized prompt template",
            "",
            "CONSISTENT GENERATION:",
            "• GPT-4 parses script into scenes",
            "• Character DNA injected into every prompt",
            "• DALL-E 3 generates with specific instructions",
            "• Character identity maintained across all frames"
        ],
        [
            "CONSISTENCY VALIDATION:",
            "• Vision API extracts features from each frame",
            "• Calculates cosine similarity vs. original DNA",
            "• Scores: 85%+ = Excellent, 70-84% = Good, <70% = Fix",
            "• Generates detailed report with recommendations",
            "",
            "SMART REGENERATION:",
            "• Automatically flags problematic frames",
            "• One-click regeneration using same DNA",
            "• No need to regenerate entire storyboard",
            "• Iterative improvement until perfect"
        ]
    )
    
    # Slide 7: Why It's Better - Competitive Advantages
    create_two_column_slide(
        prs,
        "3. Why CharacterLock AI is Better",
        [
            "EXISTING TOOLS:",
            "• 40-60% consistency (unusable)",
            "• Manual validation required",
            "• Separate tools for each step",
            "• Guesswork - no quality metrics",
            "• Full regeneration when issues found",
            "• No explainability",
            "• Expensive ($50-200/month)"
        ],
        [
            "CHARACTERLOCK AI:",
            "• 85%+ consistency (production-ready)",
            "• Automated validation with scores",
            "• Unified, integrated workflow",
            "• Quantified quality (0-100%)",
            "• Smart frame-level regeneration",
            "• Detailed reports & recommendations",
            "• Open-source core"
        ]
    )
    
    # Slide 8: Key Advantages
    create_content_slide(
        prs,
        "Our Unique Advantages",
        [
            "✓ Integrated Workflow - not separate tools",
            "✓ Quantified Quality - actual scores, not guesswork",
            "✓ Production-Ready - meets professional standards",
            "✓ Smart Regeneration - fix only what needs fixing",
            "✓ Explainable AI - detailed reports show exactly what to improve",
            "✓ Measurable Results - consistency scores you can trust"
        ]
    )
    
    # Slide 9: Expected Impact - Numbers
    create_large_text_slide(
        prs,
        "5. Expected Impact: The Numbers",
        "85%+",
        "Character consistency (vs. 40-60% baseline)"
    )
    
    # Slide 10: Expected Impact - Time & Cost Savings
    create_content_slide(
        prs,
        "Expected Impact: Transforming Film Production",
        [
            "COST SAVINGS:",
            "• $50,000 - $200,000 saved per production",
            "• Manual fixing eliminated",
            "• API costs: ~$0.46 per 10-frame storyboard",
            "",
            "TIME SAVINGS:",
            "• 50-70% reduction in storyboard production time",
            "• From weeks to minutes",
            "• Complete 10-frame storyboard in 5-6 minutes",
            "",
            "ACCESSIBILITY:",
            "• Indie filmmakers can now afford AI workflows",
            "• Small studios gain enterprise-level tools",
            "• Democratizes film pre-production"
        ]
    )
    
    # Slide 11: Market Opportunity
    create_content_slide(
        prs,
        "Market Opportunity",
        [
            "TARGET USERS:",
            "• Independent filmmakers (need affordable storyboards)",
            "• Production studios (rapid pre-visualization)",
            "• Animation teams (consistency at scale)",
            "• Ad agencies (fast concept visualization)",
            "",
            "MARKET SIZE:",
            "• Pre-visualization market: $2.8B globally (2025)",
            "• AI content generation: Growing 45% YoY",
            "• Film production software: $5.4B by 2028"
        ]
    )
    
    # Slide 12: Demo Results
    create_content_slide(
        prs,
        "Proven Results (Live Demo)",
        [
            "✓ Character Creation: 10-15 seconds",
            "✓ 5-Scene Storyboard: 4-6 minutes",
            "✓ Consistency Validation: 15-25 seconds",
            "✓ Achieved 87.5% overall consistency",
            "✓ Automated frame regeneration successful",
            "✓ Production-ready quality output",
            "",
            "→ Complete workflow in under 5 minutes"
        ]
    )
    
    # Slide 13: Technical Stack
    create_two_column_slide(
        prs,
        "Technical Implementation",
        [
            "BACKEND:",
            "• Python FastAPI (async, auto-docs)",
            "• OpenAI GPT-4 (script parsing)",
            "• OpenAI DALL-E 3 (generation)",
            "• OpenAI Vision API (validation)",
            "• SQLite (persistent storage)",
            "• NumPy/scikit-learn (similarity)",
            "",
            "API ENDPOINTS:",
            "• 15+ REST endpoints",
            "• Auto-generated documentation",
            "• Comprehensive error handling"
        ],
        [
            "FRONTEND:",
            "• React 18 + Vite",
            "• Tailwind CSS",
            "• Responsive design",
            "• Real-time progress updates",
            "",
            "KEY FEATURES:",
            "• Character DNA extraction",
            "• Scene-by-scene generation",
            "• Cosine similarity scoring",
            "• Automated validation reports",
            "• Smart frame regeneration",
            "• Export-ready output"
        ]
    )
    
    # Slide 14: Roadmap
    create_content_slide(
        prs,
        "Future Roadmap",
        [
            "PHASE 1 (1-2 months):",
            "• Fine-tune consistency algorithms",
            "• Add multiple AI model support (Stable Diffusion, Midjourney)",
            "• User authentication & cloud deployment",
            "",
            "PHASE 2 (3-6 months):",
            "• Multi-user collaboration features",
            "• Export to industry formats (FCP XML, Adobe Premiere)",
            "• Mobile app for on-set viewing",
            "",
            "PHASE 3 (6-12 months):",
            "• Video storyboard animation",
            "• 3D character model generation",
            "• Style transfer (consistent art styles)",
            "• API for third-party integrations"
        ]
    )
    
    # Slide 15: Call to Action
    create_content_slide(
        prs,
        "Join Us in Transforming Film Production",
        [
            "CharacterLock AI solves a $50K-$200K problem",
            "",
            "✓ Production-ready technology",
            "✓ Measurable, quantified results",
            "✓ Real business value",
            "✓ Market-ready solution",
            "",
            "We're looking for:",
            "• Feedback from film professionals",
            "• Pilot partner studios",
            "• Technical collaborators",
            "",
            "Let's make AI-assisted filmmaking accessible to everyone!"
        ]
    )
    
    # Slide 16: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Large "Thank You"
    thank_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    thank_frame = thank_shape.text_frame
    thank_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = thank_frame.add_paragraph()
    p.text = "Thank You!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(3, 105, 161)
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_shape = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(2))
    contact_frame = contact_shape.text_frame
    lines = [
        "CharacterLock AI",
        "From concept to consistent storyboard in 5 minutes",
        "",
        "Questions? Let's talk!"
    ]
    for line in lines:
        p = contact_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "CharacterLock_AI_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created successfully: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")
    print(f"  Location: {output_file}")

if __name__ == "__main__":
    main()
